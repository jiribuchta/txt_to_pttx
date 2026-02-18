from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
import os
import requests
from bs4 import BeautifulSoup
import flet as ft

def extract_text_from_html(song_num):
    page = requests.get(f"https://zpevnik.proscholy.cz/pisen/{song_num}")
    page.encoding = 'utf-8'
    html_content = page.text
    soup = BeautifulSoup(html_content, 'html.parser')

    lines = soup.find_all('div', class_='song-line')

    result = ""
    for line in lines:
        line_output = ""
        
        for element in line.children:
            if element.name == 'span' and 'chord' in element.get('class', []):
                chord_text = element.find('span', class_='chord-text')
                
                text = chord_text.get_text() if chord_text else ""
                
                line_output += f"{text}"
            else:
                line_output += element.get_text()
    
        if line_output.strip() == "":
            line_output = "\n"
        result += line_output + "\n"
    return result

def generate_presentation(content):
    prs = Presentation()
    prs.slide_height = Inches(9)
    prs.slide_width = Inches(16)

    font_size = 48
    verse = content.split("\n\n")

    for s in verse:
        title_slide_layout = prs.slide_layouts[6] # Blank slide layout
        slide = prs.slides.add_slide(title_slide_layout)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(16 - 2), Inches(9 - 2))
        tf = txBox.text_frame

        p = tf.add_paragraph()
        p.font.size = Pt(font_size)

        p.text = s.strip()
        p.alignment = PP_ALIGN.CENTER
    prs.save('presentation.pptx')

def split_text_to_dict(content):
    titles = ["KYRIE", "GLORIA", "CREDO", "SANCTUS", "AGNUS DEI"]
    
    lines = content.split('\n')
    structured_data = {}
    
    current_section = "ÚVOD"
    current_block = []
    
    structured_data[current_section] = []

    for line in lines:
        clean_line = line.strip()
        
        if not clean_line:
            if current_block:
                structured_data[current_section].append("\n".join(current_block))
                current_block = []
            continue

        found_title = next((t for t in titles if t in clean_line.upper()), None)
        
        if found_title:
            if current_block:
                structured_data[current_section].append("\n".join(current_block))
                current_block = []
            
            current_section = found_title
            if current_section not in structured_data:
                structured_data[current_section] = []
        else:
            current_block.append(clean_line)

    if current_block:
        structured_data[current_section].append("\n".join(current_block))
        
    if not structured_data["ÚVOD"]:
        del structured_data["ÚVOD"]
        
    return structured_data

def main():
    songs_texts = []
    while True:
        input_song_num = input("Zadejte číslo písně (nebo enter pro ukončení): ")
        if input_song_num.lower() == '':
           break
        text = extract_text_from_html(input_song_num)
        print(text)
        songs_texts.append(text)
    ordinarium = input("Chcete ordinarium? (zadej číslo nebo enter pro ukončení): ")
    text = extract_text_from_html(ordinarium)
    ordinarium_blocks = split_text_to_dict(text)
    if not all(key in ordinarium_blocks for key in ["KYRIE", "SANCTUS", "AGNUS DEI"]):
        print("Chybí některé části ordinária (KYRIE, SANCTUS, AGNUS DEI). Nelze pokračovat.")
        return

    songs_texts.insert(1, "\n".join(ordinarium_blocks["KYRIE"]))
    songs_texts.insert(4, "\n".join(ordinarium_blocks["SANCTUS"]))
    songs_texts.insert(5, "\n".join(ordinarium_blocks["AGNUS DEI"]))

    if "GLORIA" in ordinarium_blocks and "CREDO" in ordinarium_blocks:
        slavnost = input("Je slavnost? (a/n): ")
        if slavnost.lower() == 'a':
            songs_texts.insert(2, "\n".join(ordinarium_blocks["GLORIA"]))
            songs_texts.insert(4, "\n".join(ordinarium_blocks["CREDO"]))
    
    generate_presentation("\n\n".join(songs_texts))
    
if __name__ == "__main__":
    main()